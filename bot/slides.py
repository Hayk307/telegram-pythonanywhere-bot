"""Slide-deck generation for the /slides command.

The user gives a topic (and optionally how many slides); the AI returns a
structured spec (title, subtitle, and a list of slides with a heading + bullet
points); this module turns that spec into an attractive, consistently-themed
.pptx.

python-pptx is imported **lazily inside** build_deck() — never at module
import time — so the bot still boots and every other command keeps working
even when the library isn't installed on the host yet. A /slides attempt then
just fails with a clear "not available" message instead of taking the whole
worker down. This mirrors bot/fileconvert.py.

parse_deck_spec() and parse_slides_request() are pure stdlib, so the AI-output
parsing and the command parsing are testable without python-pptx installed.
"""

from __future__ import annotations

import json
import os
import re

# Bundled Unicode font so PDF export renders non-Latin text (e.g. Russian).
# fpdf2's built-in core fonts are Latin-1 only, so Cyrillic would become "?"
# without a TrueType font that has the glyphs. DejaVu Sans is freely
# redistributable and ships in bot/assets/fonts/.
_FONT_DIR = os.path.join(os.path.dirname(__file__), "assets", "fonts")
_PDF_FONT_REGULAR = os.path.join(_FONT_DIR, "DejaVuSans.ttf")
_PDF_FONT_BOLD = os.path.join(_FONT_DIR, "DejaVuSans-Bold.ttf")

VALID_FORMATS = ("pptx", "pdf")

# Upper safety ceiling. The user can ask for any number of slides up to this;
# beyond it a single AI call (bounded by Telegram's ~60s webhook window and the
# model's max output) can't reliably produce complete JSON, so we clamp and the
# handler tells the user. Bullets stay modest so slides read cleanly.
MAX_SLIDES = 50
MAX_BULLETS = 6
MAX_HEADING_LEN = 120
MAX_BULLET_LEN = 300
MAX_NOTES_LEN = 1200  # speaker notes per slide (PPTX notes pane)

# 16:9 widescreen, measured in EMU (English Metric Units; 914400 EMU = 1 inch).
_SLIDE_W = 12192000  # 13.333 in
_SLIDE_H = 6858000   # 7.5 in
_IN = 914400         # 1 inch in EMU

# A single, cohesive, modern theme so every deck looks intentional.
_PRIMARY = (0x14, 0x33, 0x55)      # deep navy — title bg + header bands
_ACCENT = (0x2F, 0x86, 0xC9)       # azure — bullet dots, rules
_HIGHLIGHT = (0xF3, 0x9C, 0x12)    # amber — top strip, title underline
_TEXT_DARK = (0x21, 0x2B, 0x36)    # body text
_TEXT_LIGHT = (0xFF, 0xFF, 0xFF)   # text on the navy background
_SUBTLE = (0xC7, 0xD3, 0xE8)       # subtitle on navy
_MUTED = (0x9A, 0x9A, 0x9A)        # footer / slide numbers
_FONT = "Calibri"


class SlideError(Exception):
    """Raised with a user-facing message when a deck can't be built."""


def _clean(value) -> str:
    """Coerce a spec value to a trimmed single-spaced string."""
    return re.sub(r"\s+", " ", str(value or "")).strip()


_FMT_WORD = r"(pdf|pptx|ppt|powerpoint|powerpoint)"


def parse_slides_request(text: str):
    """Split the /slides argument into (count | None, fmt, topic).

    - fmt is "pdf" or "pptx" (default). Set by a format word anywhere:
      "as pdf", "in pdf format", or a bare "pdf" / "pptx" token, which is
      stripped from the topic.
    - count comes from a leading number ("20 climate change") or a
      "N slides" phrase ("climate change, 20 slides"); None lets the AI choose.
    - A bare number with no topic yields (count, fmt, "").
    """
    text = _clean(text)
    fmt = "pptx"
    if not text:
        return None, fmt, ""

    m = re.search(
        r"\b(?:as|to|in)\s+(?:an?\s+)?" + _FMT_WORD + r"(?:\s+format)?\b",
        text,
        re.IGNORECASE,
    ) or re.search(r"\b" + _FMT_WORD + r"\b", text, re.IGNORECASE)
    if m:
        fmt = "pdf" if m.group(1).lower() == "pdf" else "pptx"
        text = (text[: m.start()] + " " + text[m.end() :]).strip()
        text = re.sub(r"\s+", " ", text).strip(" ,.:;-")

    m = re.match(r"^(\d{1,3})\b[\s,.:;-]*(.*)$", text)
    if m:
        return int(m.group(1)), fmt, m.group(2).strip()
    m = re.search(r"(\d{1,3})\s*slides?\b", text, re.IGNORECASE)
    if m:
        topic = (text[: m.start()] + text[m.end() :]).strip(" ,.:;-")
        return int(m.group(1)), fmt, topic
    return None, fmt, text


def parse_deck_spec(raw: str):
    """Parse the AI's response into a validated (title, subtitle, slides) tuple.

    `slides` is a list of {"heading": str, "bullets": [str, ...]}. The AI is
    asked for strict JSON, but models sometimes wrap it in ```json fences or
    add a sentence of preamble, so we extract the outermost {...} block before
    decoding. Raises SlideError (user-facing) on anything unusable.
    """
    text = (raw or "").strip()
    if text.startswith("```"):
        text = re.sub(r"^```[a-zA-Z]*\s*", "", text)
        text = re.sub(r"\s*```$", "", text).strip()
    start, end = text.find("{"), text.rfind("}")
    if start == -1 or end == -1 or end <= start:
        raise SlideError("The AI didn't return slide data. Try rephrasing the topic.")
    try:
        data = json.loads(text[start : end + 1])
    except json.JSONDecodeError:
        raise SlideError(
            "The AI's slide data was incomplete — try again, or ask for fewer slides."
        )

    title = _clean(data.get("title")) or "Presentation"
    subtitle = _clean(data.get("subtitle"))

    slides = []
    for item in (data.get("slides") or [])[:MAX_SLIDES]:
        if not isinstance(item, dict):
            continue
        heading = _clean(item.get("heading") or item.get("title"))[:MAX_HEADING_LEN]
        bullets = []
        for bullet in (item.get("bullets") or [])[:MAX_BULLETS]:
            bullet = _clean(bullet)[:MAX_BULLET_LEN]
            if bullet:
                bullets.append(bullet)
        notes = _clean(item.get("notes"))[:MAX_NOTES_LEN]
        if heading or bullets:
            slides.append(
                {"heading": heading or title, "bullets": bullets, "notes": notes}
            )

    if not slides:
        raise SlideError("The AI returned no usable slides. Try a clearer topic.")
    return title, subtitle, slides


# Namespaces for the raw animation XML we inject (python-pptx exposes no API).
_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _add_fade_transition(slide) -> None:
    """Add a fade transition when advancing TO this slide."""
    from pptx.oxml import parse_xml

    xml = (
        f'<p:transition xmlns:p="{_P_NS}" xmlns:a="{_A_NS}" spd="med">'
        f"<p:fade/></p:transition>"
    )
    # Schema order in <p:sld> is cSld, clrMapOvr, transition, timing — both
    # transition and timing come after the existing children, so append works.
    slide._element.append(parse_xml(xml))


def _add_bullet_build(slide, shape_id: int, n_paragraphs: int) -> None:
    """Inject a per-paragraph 'fade in on click' entrance animation.

    Builds the <p:timing> tree PowerPoint uses for a "Fade / By paragraph"
    entrance: each paragraph of the bullet text box (identified by shape_id)
    reveals on its own click. This is hand-written OOXML because python-pptx
    has no animation API; malformed timing would make PowerPoint drop the
    animation (or refuse the file), so the structure mirrors what PowerPoint
    itself emits and is covered by a round-trip test.
    """
    from pptx.oxml import parse_xml

    # Unique cTn ids across the whole tree; 1 and 2 are the root + mainSeq.
    counter = [3]

    def next_id():
        val = counter[0]
        counter[0] += 1
        return val

    effects = []
    for i in range(n_paragraphs):
        a, b, c, d, e = (next_id() for _ in range(5))
        tgt = (
            f'<p:spTgt spid="{shape_id}"><p:txEl>'
            f'<p:pRg st="{i}" end="{i}"/></p:txEl></p:spTgt>'
        )
        effects.append(
            f"<p:par><p:cTn id=\"{a}\" fill=\"hold\">"
            f'<p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>'
            f"<p:childTnLst><p:par><p:cTn id=\"{b}\" fill=\"hold\">"
            f'<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
            f"<p:childTnLst><p:par><p:cTn id=\"{c}\" presetID=\"10\" "
            f'presetClass="entr" presetSubtype="0" fill="hold" grpId="0" '
            f'nodeType="clickEffect"><p:stCondLst><p:cond delay="0"/>'
            f"</p:stCondLst><p:childTnLst>"
            f"<p:set><p:cBhvr><p:cTn id=\"{d}\" dur=\"1\" fill=\"hold\">"
            f'<p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>'
            f"<p:tgtEl>{tgt}</p:tgtEl>"
            f"<p:attrNameLst><p:attrName>style.visibility</p:attrName>"
            f"</p:attrNameLst></p:cBhvr>"
            f'<p:to><p:strVal val="visible"/></p:to></p:set>'
            f'<p:animEffect transition="in" filter="fade"><p:cBhvr>'
            f"<p:cTn id=\"{e}\" dur=\"500\"/>"
            f"<p:tgtEl>{tgt}</p:tgtEl></p:cBhvr></p:animEffect>"
            f"</p:childTnLst></p:cTn></p:par></p:childTnLst>"
            f"</p:cTn></p:par></p:childTnLst></p:cTn></p:par>"
        )

    xml = (
        f'<p:timing xmlns:p="{_P_NS}" xmlns:a="{_A_NS}"><p:tnLst><p:par>'
        f'<p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">'
        f'<p:childTnLst><p:seq concurrent="1" nextAc="seek">'
        f'<p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>'
        f"{''.join(effects)}"
        f"</p:childTnLst></p:cTn>"
        f'<p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl>'
        f"<p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>"
        f'<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl>'
        f"<p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>"
        f"</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst>"
        f'<p:bldLst><p:bldP spid="{shape_id}" grpId="0" build="p"/></p:bldLst>'
        f"</p:timing>"
    )
    slide._element.append(parse_xml(xml))


def build_deck(title: str, subtitle: str, slides, output_path: str) -> str:
    """Render the spec into an attractive, themed .pptx at output_path.

    Content slides get speaker notes and a per-bullet fade-in entrance
    animation. Raises SlideError if python-pptx isn't installed.
    """
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        from pptx.util import Emu, Pt
    except ImportError:
        raise SlideError("Slide generation isn't available (missing python-pptx).")

    def rgb(triple):
        return RGBColor(*triple)

    def rect(slide, left, top, width, height, color):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Emu(left), Emu(top), Emu(width), Emu(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(color)
        shape.line.fill.background()
        shape.shadow.inherit = False
        return shape

    def add_run(paragraph, text, *, size, color, bold=False):
        run = paragraph.add_run()
        run.text = text
        run.font.name = _FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
        return run

    prs = Presentation()
    prs.slide_width = Emu(_SLIDE_W)
    prs.slide_height = Emu(_SLIDE_H)
    blank = prs.slide_layouts[6]  # fully blank — we place every element

    total = len(slides)
    footer_title = (title[:57] + "…") if len(title) > 58 else title

    # ── Title slide ─────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    rect(slide, 0, 0, _SLIDE_W, _SLIDE_H, _PRIMARY)          # full navy field
    rect(slide, 0, 0, _SLIDE_W, int(_IN * 0.18), _HIGHLIGHT)  # amber top strip
    rect(slide, 0, _SLIDE_H - int(_IN * 0.18), _SLIDE_W, int(_IN * 0.18), _ACCENT)

    box = slide.shapes.add_textbox(
        Emu(_IN), Emu(int(_SLIDE_H * 0.33)), Emu(_SLIDE_W - 2 * _IN), Emu(int(_IN * 2.6))
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tp = tf.paragraphs[0]
    tp.alignment = PP_ALIGN.CENTER
    add_run(tp, title, size=42, color=_TEXT_LIGHT, bold=True)
    if subtitle:
        sp = tf.add_paragraph()
        sp.alignment = PP_ALIGN.CENTER
        sp.space_before = Pt(16)
        add_run(sp, subtitle, size=22, color=_SUBTLE)
    # Short amber underline centered beneath the title block.
    rect(
        slide,
        int(_SLIDE_W / 2 - _IN),
        int(_SLIDE_H * 0.33) + int(_IN * 2.6) + int(_IN * 0.1),
        2 * _IN,
        int(_IN * 0.06),
        _HIGHLIGHT,
    )

    # ── Content slides ────────────────────────────────────────────────────
    margin = int(_IN * 0.8)
    content_w = _SLIDE_W - 2 * margin
    band_h = int(_IN * 1.25)
    for index, spec in enumerate(slides, start=1):
        slide = prs.slides.add_slide(blank)

        # Header band with a thin amber accent strip on top.
        rect(slide, 0, 0, _SLIDE_W, band_h, _PRIMARY)
        rect(slide, 0, 0, _SLIDE_W, int(_IN * 0.12), _HIGHLIGHT)

        # Heading inside the band (left), vertically centered.
        head_box = slide.shapes.add_textbox(
            Emu(margin), Emu(0), Emu(content_w - int(_IN * 1.1)), Emu(band_h)
        )
        htf = head_box.text_frame
        htf.word_wrap = True
        htf.vertical_anchor = MSO_ANCHOR.MIDDLE
        add_run(htf.paragraphs[0], spec["heading"], size=28, color=_TEXT_LIGHT, bold=True)

        # Big faint slide number on the right of the band.
        num_box = slide.shapes.add_textbox(
            Emu(_SLIDE_W - int(_IN * 1.6)), Emu(0), Emu(int(_IN * 1.2)), Emu(band_h)
        )
        ntf = num_box.text_frame
        ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
        np = ntf.paragraphs[0]
        np.alignment = PP_ALIGN.RIGHT
        add_run(np, str(index), size=30, color=_ACCENT, bold=True)

        # Bullets — a colored dot run + dark text run per point. Font eases
        # down a touch on dense slides so they still fit cleanly.
        bullets = spec["bullets"]
        size = 20 if len(bullets) <= 5 else 18
        body_box = slide.shapes.add_textbox(
            Emu(margin),
            Emu(band_h + int(_IN * 0.35)),
            Emu(content_w),
            Emu(_SLIDE_H - band_h - int(_IN * 0.9)),
        )
        btf = body_box.text_frame
        btf.word_wrap = True
        for i, bullet in enumerate(bullets):
            para = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
            para.space_after = Pt(14)
            add_run(para, "●  ", size=size, color=_ACCENT, bold=True)
            add_run(para, bullet, size=size, color=_TEXT_DARK)

        # Speaker notes (elaboration the AI wrote) go in the notes pane.
        if spec.get("notes"):
            slide.notes_slide.notes_text_frame.text = spec["notes"]

        # Entrance animation: fade each bullet in one-by-one on click, plus a
        # slide fade transition. python-pptx has no animation API, so we inject
        # the timing XML directly (see _add_bullet_build / _add_fade_transition).
        if bullets:
            _add_bullet_build(slide, body_box.shape_id, len(bullets))
        _add_fade_transition(slide)

        # Footer: thin accent line, deck title (left), n / total (right).
        rect(
            slide,
            margin,
            _SLIDE_H - int(_IN * 0.55),
            content_w,
            int(_IN * 0.015),
            _ACCENT,
        )
        foot = slide.shapes.add_textbox(
            Emu(margin), Emu(_SLIDE_H - int(_IN * 0.5)), Emu(content_w), Emu(int(_IN * 0.4))
        )
        fp = foot.text_frame.paragraphs[0]
        add_run(fp, footer_title, size=10, color=_MUTED)
        fnum = slide.shapes.add_textbox(
            Emu(margin), Emu(_SLIDE_H - int(_IN * 0.5)), Emu(content_w), Emu(int(_IN * 0.4))
        )
        fnp = fnum.text_frame.paragraphs[0]
        fnp.alignment = PP_ALIGN.RIGHT
        add_run(fnp, f"{index} / {total}", size=10, color=_MUTED)

    prs.save(output_path)
    return output_path


def build_pdf(title: str, subtitle: str, slides, output_path: str) -> str:
    """Render the spec into a themed PDF at output_path (same look as the .pptx).

    Uses fpdf2 (already a dependency, imported lazily) with a bundled DejaVu
    Sans TTF so non-Latin text (Russian) renders instead of turning into "?".
    Raises SlideError if fpdf2 or the bundled font is unavailable.
    """
    try:
        from fpdf import FPDF
    except ImportError:
        raise SlideError("PDF export isn't available (missing fpdf2).")
    if not os.path.exists(_PDF_FONT_REGULAR):
        raise SlideError("PDF export isn't available (bundled font missing).")

    # 16:9 page in points (960 x 540 pt == 13.333 x 7.5 in), matching the .pptx.
    W, H = 960.0, 540.0
    MARGIN = 60.0
    BAND_H = 90.0
    CONTENT_W = W - 2 * MARGIN

    pdf = FPDF(orientation="L", unit="pt", format=(W, H))
    pdf.set_auto_page_break(False)  # one slide == one page; we place everything
    pdf.set_margins(MARGIN, MARGIN, MARGIN)
    pdf.add_font("DejaVu", "", _PDF_FONT_REGULAR)
    if os.path.exists(_PDF_FONT_BOLD):
        pdf.add_font("DejaVu", "B", _PDF_FONT_BOLD)
    else:
        # Fall back to the regular face for "bold" so set_font never fails.
        pdf.add_font("DejaVu", "B", _PDF_FONT_REGULAR)

    def fill(rgb):
        pdf.set_fill_color(*rgb)

    def text_color(rgb):
        pdf.set_text_color(*rgb)

    def bar(x, y, w, h, rgb):
        fill(rgb)
        pdf.rect(x, y, w, h, style="F")

    total = len(slides)
    footer_title = (title[:57] + "…") if len(title) > 58 else title

    # ── Title page ───────────────────────────────────────────────────────────
    pdf.add_page()
    bar(0, 0, W, H, _PRIMARY)          # full navy field
    bar(0, 0, W, 10, _HIGHLIGHT)       # amber top strip
    bar(0, H - 10, W, 10, _ACCENT)     # azure bottom strip

    pdf.set_font("DejaVu", "B", 34)
    text_color(_TEXT_LIGHT)
    pdf.set_xy(MARGIN, H * 0.34)
    pdf.multi_cell(CONTENT_W, 42, title, align="C")
    if subtitle:
        pdf.set_font("DejaVu", "", 18)
        text_color(_SUBTLE)
        pdf.set_xy(MARGIN, pdf.get_y() + 12)
        pdf.multi_cell(CONTENT_W, 24, subtitle, align="C")
    underline_y = pdf.get_y() + 14
    bar(W / 2 - 72, underline_y, 144, 5, _HIGHLIGHT)  # centered amber underline

    # ── Content pages ─────────────────────────────────────────────────────────
    for index, spec in enumerate(slides, start=1):
        pdf.add_page()
        bar(0, 0, W, BAND_H, _PRIMARY)  # header band
        bar(0, 0, W, 8, _HIGHLIGHT)     # amber accent strip

        # Heading (white, left), vertically centered in the band.
        pdf.set_font("DejaVu", "B", 22)
        text_color(_TEXT_LIGHT)
        pdf.set_xy(MARGIN, BAND_H / 2 - 16)
        pdf.multi_cell(CONTENT_W - 90, 26, spec["heading"], align="L")

        # Slide number (azure) on the right of the band.
        pdf.set_font("DejaVu", "B", 24)
        text_color(_ACCENT)
        pdf.set_xy(W - MARGIN - 90, BAND_H / 2 - 16)
        pdf.cell(90, 28, str(index), align="R")

        # Bullets: azure dot + dark wrapping text.
        bullets = spec["bullets"]
        size = 15 if len(bullets) <= 5 else 13
        line_h = size * 1.4
        dot_w = size * 1.4
        pdf.set_y(BAND_H + 28)
        for bullet in bullets:
            y0 = pdf.get_y()
            pdf.set_font("DejaVu", "B", size)
            text_color(_ACCENT)
            pdf.set_xy(MARGIN, y0)
            pdf.cell(dot_w, line_h, "•")
            pdf.set_font("DejaVu", "", size)
            text_color(_TEXT_DARK)
            pdf.set_xy(MARGIN + dot_w, y0)
            pdf.multi_cell(
                CONTENT_W - dot_w, line_h, bullet, align="L", new_x="LMARGIN", new_y="NEXT"
            )
            pdf.set_y(pdf.get_y() + 8)

        # Footer: thin accent rule, deck title (left), n / total (right).
        bar(MARGIN, H - 38, CONTENT_W, 1.2, _ACCENT)
        pdf.set_font("DejaVu", "", 9)
        text_color(_MUTED)
        pdf.set_xy(MARGIN, H - 32)
        pdf.cell(CONTENT_W / 2, 12, footer_title, align="L")
        pdf.set_xy(MARGIN + CONTENT_W / 2, H - 32)
        pdf.cell(CONTENT_W / 2, 12, f"{index} / {total}", align="R")

    pdf.output(output_path)
    return output_path
