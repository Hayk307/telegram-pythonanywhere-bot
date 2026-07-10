"""Tests for bot/slides.py.

The parsing/validation tests need no third-party libraries. The real
deck-building test is guarded with pytest.importorskip so it exercises the
actual python-pptx renderer when the optional dep is installed (as in CI via
requirements.txt) and skips otherwise.
"""

import json

import pytest

from bot import slides


# ── parse_slides_request: command parsing, no optional deps ─────────────────


def test_request_no_count():
    assert slides.parse_slides_request("история интернета") == (
        None,
        "pptx",
        "история интернета",
    )


def test_request_leading_count():
    assert slides.parse_slides_request("15 основы фотосинтеза") == (
        15,
        "pptx",
        "основы фотосинтеза",
    )


def test_request_leading_count_with_separator():
    assert slides.parse_slides_request("20, климат") == (20, "pptx", "климат")


def test_request_n_slides_phrase():
    count, fmt, topic = slides.parse_slides_request("как работает блокчейн, 20 slides")
    assert count == 20
    assert fmt == "pptx"
    assert topic == "как работает блокчейн"


def test_request_bare_number_no_topic():
    assert slides.parse_slides_request("12") == (12, "pptx", "")


def test_request_empty():
    assert slides.parse_slides_request("") == (None, "pptx", "")
    assert slides.parse_slides_request("   ") == (None, "pptx", "")


def test_request_number_inside_topic_is_not_a_count():
    # No leading number and no "N slides" phrase -> the digits stay in the topic.
    assert slides.parse_slides_request("top 10 programming languages") == (
        None,
        "pptx",
        "top 10 programming languages",
    )


# ── format detection ────────────────────────────────────────────────────────


def test_request_pdf_trailing_token():
    count, fmt, topic = slides.parse_slides_request("12 фотосинтез pdf")
    assert (count, fmt, topic) == (12, "pdf", "фотосинтез")


def test_request_pdf_as_phrase():
    count, fmt, topic = slides.parse_slides_request("климат as pdf")
    assert count is None
    assert fmt == "pdf"
    assert topic == "климат"


def test_request_pdf_in_format_phrase():
    _, fmt, topic = slides.parse_slides_request("15 блокчейн in pdf format")
    assert fmt == "pdf"
    assert topic == "блокчейн"


def test_request_pptx_explicit():
    _, fmt, topic = slides.parse_slides_request("история powerpoint")
    assert fmt == "pptx"
    assert topic == "история"


def test_request_default_format_is_pptx():
    _, fmt, _ = slides.parse_slides_request("любая тема")
    assert fmt == "pptx"


# ── parse_deck_spec: pure logic, no optional deps ───────────────────────────


def _spec(**kw):
    base = {
        "title": "Интернет",
        "subtitle": "краткая история",
        "slides": [
            {"heading": "Начало", "bullets": ["ARPANET", "1969"]},
            {"heading": "Веб", "bullets": ["1991", "HTTP"]},
        ],
    }
    base.update(kw)
    return json.dumps(base, ensure_ascii=False)


def test_parse_basic():
    title, subtitle, deck = slides.parse_deck_spec(_spec())
    assert title == "Интернет"
    assert subtitle == "краткая история"
    assert len(deck) == 2
    assert deck[0]["heading"] == "Начало"
    assert deck[0]["bullets"] == ["ARPANET", "1969"]


def test_parse_strips_code_fences():
    raw = "```json\n" + _spec() + "\n```"
    title, _, deck = slides.parse_deck_spec(raw)
    assert title == "Интернет"
    assert len(deck) == 2


def test_parse_tolerates_prose_around_json():
    raw = "Sure! Here is your deck:\n" + _spec() + "\nHope that helps."
    _, _, deck = slides.parse_deck_spec(raw)
    assert len(deck) == 2


def test_parse_defaults_title_and_optional_subtitle():
    raw = json.dumps({"slides": [{"heading": "H", "bullets": ["b"]}]})
    title, subtitle, deck = slides.parse_deck_spec(raw)
    assert title == "Presentation"
    assert subtitle == ""
    assert len(deck) == 1


def test_parse_clamps_slide_and_bullet_counts():
    many = {
        "title": "T",
        "slides": [
            {"heading": f"H{i}", "bullets": [f"b{j}" for j in range(20)]}
            for i in range(slides.MAX_SLIDES + 15)
        ],
    }
    _, _, deck = slides.parse_deck_spec(json.dumps(many))
    assert len(deck) == slides.MAX_SLIDES
    assert all(len(s["bullets"]) <= slides.MAX_BULLETS for s in deck)


def test_parse_accepts_title_key_for_slide_heading():
    raw = json.dumps({"title": "T", "slides": [{"title": "Alt", "bullets": ["x"]}]})
    _, _, deck = slides.parse_deck_spec(raw)
    assert deck[0]["heading"] == "Alt"


def test_parse_skips_empty_slides():
    raw = json.dumps(
        {
            "title": "T",
            "slides": [
                {"heading": "", "bullets": []},
                {"heading": "Real", "bullets": ["ok"]},
            ],
        }
    )
    _, _, deck = slides.parse_deck_spec(raw)
    assert len(deck) == 1
    assert deck[0]["heading"] == "Real"


def test_parse_no_json_raises():
    with pytest.raises(slides.SlideError):
        slides.parse_deck_spec("I could not make slides, sorry.")


def test_parse_malformed_json_raises():
    with pytest.raises(slides.SlideError):
        slides.parse_deck_spec('{"title": "T", "slides": [ oops }')


def test_parse_no_usable_slides_raises():
    with pytest.raises(slides.SlideError):
        slides.parse_deck_spec(json.dumps({"title": "T", "slides": []}))


# ── build_deck: real render, needs python-pptx ──────────────────────────────


def test_build_deck_creates_pptx(tmp_path):
    pytest.importorskip("pptx")
    out = tmp_path / "deck.pptx"
    title, subtitle, deck = slides.parse_deck_spec(_spec())
    result = slides.build_deck(title, subtitle, deck, str(out))
    assert result == str(out)
    assert out.exists() and out.stat().st_size > 0

    from pptx import Presentation

    prs = Presentation(str(out))
    # 1 title slide + 2 content slides.
    assert len(prs.slides) == 3
    # The title text appears somewhere on the first slide.
    texts = [
        shape.text_frame.text
        for shape in prs.slides[0].shapes
        if shape.has_text_frame
    ]
    assert any("Интернет" in t for t in texts)


def test_build_pdf_creates_pdf(tmp_path):
    pytest.importorskip("fpdf")
    out = tmp_path / "deck.pdf"
    title, subtitle, deck = slides.parse_deck_spec(_spec())
    result = slides.build_pdf(title, subtitle, deck, str(out))
    assert result == str(out)
    assert out.exists()
    data = out.read_bytes()
    assert data[:5] == b"%PDF-"  # valid PDF signature
    assert len(data) > 1000

